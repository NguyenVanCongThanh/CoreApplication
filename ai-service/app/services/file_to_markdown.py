"""
ai-service/app/services/file_to_markdown.py

One canonical pipeline that turns any uploaded artefact (PDF, DOCX,
PPTX, XLSX, image, YouTube URL) into a *normalised Markdown document*
with:

  * Real Markdown headings (so MarkdownChunker breadcrumbs work)
  * Tables in Markdown table syntax (so the table-prefix logic kicks in)
  * Embedded images uploaded to MinIO and inlined as ![alt](url)
  * VLM-based OCR fallback for scanned PDFs (no Tesseract dependency)

This is the foundation of both:
  - Phase B: micro-lesson generation (needs structured Markdown to feed
    the splitter LLM and to keep figure references intact).
  - Phase A: smarter auto-indexing (the unified Markdown is fed back
    into the existing MarkdownChunker which already produces the
    best-quality semantic chunks in the codebase).

Returned object includes both the Markdown body and the list of
ExtractedImage records — the caller can decide whether to call the
VLM on the URLs (auto-index does it during chunking) or pass them to
the lesson splitter (micro-lesson generator).
"""
from __future__ import annotations

import io
import logging
import re
import uuid
from dataclasses import dataclass, field
from typing import Optional

from app.core.config import get_settings
from app.services.image_extractor import (
    ExtractedImage,
    extract_docx_images,
    extract_pdf_images,
    render_pptx_slides,
)
from app.services.minio_storage import upload_bytes

logger = logging.getLogger(__name__)
settings = get_settings()


# Heuristic: a PDF page is "scanned" (no real text layer) when extracted
# text is shorter than this threshold. Tuned on Vietnamese textbook scans.
_OCR_TEXT_THRESHOLD_CHARS = 40

# Cap how many scanned pages we'll send through VLM OCR per document so
# a single 500-page scan doesn't blow the Groq budget.
_OCR_MAX_PAGES = 60

# Render scanned pages at this DPI for the VLM. 200 keeps tables readable
# while staying under 4 MB per image.
_OCR_RENDER_DPI = 200


@dataclass
class ConvertedDocument:
    markdown: str
    images: list[ExtractedImage] = field(default_factory=list)
    page_count: int = 0
    ocr_pages: int = 0
    language_hint: str = "vi"


# ── Public entry point ────────────────────────────────────────────────────────

async def convert_to_markdown(
    file_bytes: bytes,
    file_type: str,
    storage_prefix: str,
    language: str = "vi",
) -> ConvertedDocument:
    """
    Dispatch on file_type ('pdf' | 'docx' | 'pptx' | 'xlsx' | 'text' | 'image').
    `storage_prefix` is the MinIO path prefix used for any extracted assets.
    """
    file_type = (file_type or "").lower()
    if file_type == "pdf":
        return await _pdf_to_markdown(file_bytes, storage_prefix, language)
    if file_type == "docx":
        return await _docx_to_markdown(file_bytes, storage_prefix)
    if file_type == "pptx":
        return await _pptx_to_markdown(file_bytes, storage_prefix, language)
    if file_type == "xlsx":
        return await _xlsx_to_markdown(file_bytes)
    if file_type == "image":
        return await _image_to_markdown(file_bytes, storage_prefix, language)
    # Plain text / markdown — pass through (caller may already have markdown)
    text = file_bytes.decode("utf-8", errors="replace")
    return ConvertedDocument(markdown=text)


# ── PDF ───────────────────────────────────────────────────────────────────────

async def _pdf_to_markdown(
    pdf_bytes: bytes,
    storage_prefix: str,
    language: str,
) -> ConvertedDocument:
    try:
        import pymupdf
    except ImportError:
        logger.error("PyMuPDF missing — falling back to plain pypdf extract")
        return await _pdf_to_markdown_pypdf(pdf_bytes)

    images = await extract_pdf_images(pdf_bytes, storage_prefix)
    images_by_page: dict[int, list[ExtractedImage]] = {}
    for img in images:
        if img.page_number is None:
            continue
        images_by_page.setdefault(img.page_number, []).append(img)

    doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
    total_pages = len(doc)
    out_parts: list[str] = []
    ocr_pages = 0

    try:
        for page_idx in range(total_pages):
            page = doc[page_idx]
            page_no = page_idx + 1
            page_md = _pdf_page_to_markdown(page)

            if len(page_md.strip()) < _OCR_TEXT_THRESHOLD_CHARS and ocr_pages < _OCR_MAX_PAGES:
                # Scanned page — render and OCR via VLM
                ocr_md = await _vlm_ocr_pdf_page(
                    page=page,
                    page_no=page_no,
                    storage_prefix=storage_prefix,
                    language=language,
                )
                if ocr_md:
                    page_md = ocr_md
                    ocr_pages += 1

            if not page_md.strip():
                continue

            out_parts.append(f"\n\n## Trang {page_no}\n\n")
            out_parts.append(page_md.rstrip())

            for img in images_by_page.get(page_no, []):
                alt = img.caption_hint or f"Hình minh họa trang {page_no}"
                out_parts.append(f"\n\n![{alt}]({img.url})\n")
    finally:
        doc.close()

    return ConvertedDocument(
        markdown="".join(out_parts).strip(),
        images=images,
        page_count=total_pages,
        ocr_pages=ocr_pages,
        language_hint=language,
    )


def _pdf_page_to_markdown(page) -> str:
    """
    Convert a PyMuPDF page to Markdown using font-size heuristics:
    larger blocks become headings (## or ###), normal text becomes
    paragraphs. Keeps inline reading order via blocks=True.
    """
    try:
        info = page.get_text("dict")
    except Exception:
        return page.get_text("text") or ""

    blocks = info.get("blocks", [])
    if not blocks:
        return page.get_text("text") or ""

    sizes: list[float] = []
    for b in blocks:
        for line in b.get("lines", []):
            for span in line.get("spans", []):
                sz = span.get("size") or 0
                if sz:
                    sizes.append(sz)
    if not sizes:
        return page.get_text("text") or ""

    body_size = sorted(sizes)[len(sizes) // 2]
    h2_threshold = body_size * 1.4
    h3_threshold = body_size * 1.15

    lines_out: list[str] = []
    for b in blocks:
        if b.get("type") != 0:
            continue
        for line in b.get("lines", []):
            text_parts: list[str] = []
            max_size = 0.0
            for span in line.get("spans", []):
                t = span.get("text", "")
                if t:
                    text_parts.append(t)
                max_size = max(max_size, span.get("size") or 0)
            text = " ".join(text_parts).strip()
            if not text:
                continue
            if max_size >= h2_threshold and len(text) < 120:
                lines_out.append(f"\n### {text}\n")
            elif max_size >= h3_threshold and len(text) < 140:
                lines_out.append(f"\n**{text}**\n")
            else:
                lines_out.append(text)

    md = "\n".join(lines_out)
    md = re.sub(r"\n{3,}", "\n\n", md)
    return md


async def _vlm_ocr_pdf_page(
    page,
    page_no: int,
    storage_prefix: str,
    language: str,
) -> Optional[str]:
    """
    Render the page to PNG and ask the VLM to transcribe it as Markdown.
    Image bytes are NOT uploaded — they're sent inline to the VLM.
    """
    try:
        import pymupdf

        zoom = _OCR_RENDER_DPI / 72.0
        pix = page.get_pixmap(matrix=pymupdf.Matrix(zoom, zoom), alpha=False)
        png_bytes = pix.tobytes("png")
        pix = None
    except Exception as exc:
        logger.warning("OCR render failed page=%d: %s", page_no, exc)
        return None

    try:
        from groq import AsyncGroq
        import base64

        if not settings.groq_api_key:
            return None

        prompt = (
            "Bạn nhận được ảnh chụp một trang tài liệu học thuật bị scan (không có "
            "lớp text). Hãy trích xuất TOÀN BỘ nội dung trang dưới dạng Markdown:\n"
            "- Giữ nguyên các tiêu đề lớn / nhỏ (dùng ##, ###).\n"
            "- Giữ danh sách (ordered/unordered) đúng định dạng.\n"
            "- Bảng phải dùng cú pháp bảng Markdown (| col | col |).\n"
            "- KHÔNG thêm câu mở đầu kiểu 'Đây là nội dung trang…'. Bắt đầu thẳng vào nội dung.\n"
            "- Nếu trang trống, trả về chuỗi rỗng."
            if language == "vi"
            else
            "You are given a page from a scanned academic document with no text layer. "
            "Transcribe the full page as Markdown:\n"
            "- Keep heading hierarchy (##, ###).\n"
            "- Preserve ordered/unordered lists.\n"
            "- Tables must use Markdown table syntax.\n"
            "- Do not add any preamble like 'This page contains…'. Start with the content.\n"
            "- If the page is blank, return an empty string."
        )

        b64 = base64.b64encode(png_bytes).decode("utf-8")
        client = AsyncGroq(api_key=settings.groq_api_key)
        response = await client.chat.completions.create(
            model=settings.vlm_model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                        {"type": "text", "text": prompt},
                    ],
                },
            ],
            temperature=0.0,
            max_tokens=2048,
        )
        return (response.choices[0].message.content or "").strip()
    except Exception as exc:
        logger.warning("VLM OCR failed page=%d: %s", page_no, exc)
        return None


async def _pdf_to_markdown_pypdf(pdf_bytes: bytes) -> ConvertedDocument:
    try:
        import pypdf
    except ImportError:
        return ConvertedDocument(markdown="")

    reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(f"\n\n## Trang {i + 1}\n\n{text}")
    return ConvertedDocument(markdown="".join(parts).strip(), page_count=len(reader.pages))


# ── DOCX ──────────────────────────────────────────────────────────────────────

async def _docx_to_markdown(docx_bytes: bytes, storage_prefix: str) -> ConvertedDocument:
    """
    Use mammoth to convert .docx → Markdown. Mammoth preserves heading
    levels, bold/italic, lists and tables better than python-docx.
    Images embedded in the document are also extracted to MinIO.
    """
    images = await extract_docx_images(docx_bytes, storage_prefix)
    md = ""

    try:
        import mammoth

        # Custom image converter: mammoth gives us each image; we already
        # uploaded them above, so here we just emit a placeholder that the
        # caller can splice. We rely on key order matching extraction order.
        seen_idx = {"i": 0}
        url_pool = [img.url for img in images]

        def _img_handler(image):
            idx = seen_idx["i"]
            seen_idx["i"] += 1
            url = url_pool[idx] if idx < len(url_pool) else ""
            alt = (image.alt_text or "").strip() or f"Hình {idx + 1}"
            return {"src": url, "alt": alt}

        result = mammoth.convert_to_markdown(
            io.BytesIO(docx_bytes),
            convert_image=mammoth.images.img_element(_img_handler),
        )
        md = result.value or ""
    except ImportError:
        logger.warning("mammoth not installed — using python-docx fallback")
        md = _docx_to_markdown_fallback(docx_bytes)
    except Exception as exc:
        logger.error("mammoth conversion failed: %s — falling back", exc, exc_info=True)
        md = _docx_to_markdown_fallback(docx_bytes)

    return ConvertedDocument(markdown=md.strip(), images=images)


def _docx_to_markdown_fallback(docx_bytes: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""
    doc = Document(io.BytesIO(docx_bytes))
    out: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if "heading 1" in style:
            out.append(f"\n# {text}\n")
        elif "heading 2" in style:
            out.append(f"\n## {text}\n")
        elif "heading 3" in style:
            out.append(f"\n### {text}\n")
        else:
            out.append(text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [c.text.strip().replace("|", "\\|") for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if len(rows) >= 1:
            cols = len(table.rows[0].cells)
            out.append("\n" + rows[0])
            out.append("|" + "|".join(["---"] * cols) + "|")
            for r in rows[1:]:
                out.append(r)
    return "\n".join(out)


# ── PPTX ──────────────────────────────────────────────────────────────────────

async def _pptx_to_markdown(pptx_bytes: bytes, storage_prefix: str, language: str) -> ConvertedDocument:
    """
    Convert each slide to a `## Slide N: <title>` section. Inline images
    embedded in the deck are extracted to MinIO.
    """
    images = await render_pptx_slides(pptx_bytes, storage_prefix)
    parts: list[str] = []
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        for i, slide in enumerate(prs.slides):
            slide_no = i + 1
            title = _pptx_slide_title(slide) or f"Slide {slide_no}"
            parts.append(f"\n\n## Slide {slide_no}: {title}\n")

            body_lines: list[str] = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                txt = (shape.text or "").strip()
                if not txt or txt == title:
                    continue
                for line in txt.splitlines():
                    line = line.strip()
                    if line:
                        body_lines.append(f"- {line}")
            if body_lines:
                parts.append("\n".join(body_lines))

            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(f"\n\n> **Ghi chú giảng viên:** {notes}\n")
    except Exception as exc:
        logger.error("PPTX → Markdown failed: %s", exc, exc_info=True)

    md = "\n".join(parts).strip()
    if images and "![" not in md:
        # Append image gallery at the end so they're not lost — the splitter
        # LLM can attach them to the right lesson based on caption_hint.
        md += "\n\n## Hình ảnh đính kèm\n"
        for img in images:
            md += f"\n![Hình {img.order_in_page}]({img.url})\n"
    return ConvertedDocument(markdown=md, images=images)


def _pptx_slide_title(slide) -> str:
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip().splitlines()[0][:120]
    except Exception:
        pass
    return ""


# ── XLSX ──────────────────────────────────────────────────────────────────────

async def _xlsx_to_markdown(xlsx_bytes: bytes) -> ConvertedDocument:
    """Each sheet → its own `## Sheet: <name>` section with a real Markdown table."""
    parts: list[str] = []
    try:
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True, read_only=True)
        for sheet in wb.worksheets:
            parts.append(f"\n\n## Sheet: {sheet.title}\n")
            rows = list(sheet.iter_rows(values_only=True))
            rows = [r for r in rows if any(c is not None and str(c).strip() != "" for c in r)]
            if not rows:
                continue
            header = rows[0]
            cols = len(header)
            header_cells = [_clean_cell(c) for c in header]
            parts.append("| " + " | ".join(header_cells) + " |")
            parts.append("|" + "|".join(["---"] * cols) + "|")
            for r in rows[1: 1 + 1000]:  # cap each sheet at 1k data rows
                cells = [_clean_cell(c) for c in r]
                while len(cells) < cols:
                    cells.append("")
                parts.append("| " + " | ".join(cells[:cols]) + " |")
            if len(rows) > 1001:
                parts.append(f"\n_…đã rút gọn còn 1000/{len(rows) - 1} dòng dữ liệu_\n")
    except Exception as exc:
        logger.error("XLSX → Markdown failed: %s", exc, exc_info=True)

    return ConvertedDocument(markdown="\n".join(parts).strip())


def _clean_cell(value) -> str:
    if value is None:
        return ""
    s = str(value).strip()
    return s.replace("|", "\\|").replace("\n", " ")


# ── Standalone image ──────────────────────────────────────────────────────────

async def _image_to_markdown(
    image_bytes: bytes,
    storage_prefix: str,
    language: str,
) -> ConvertedDocument:
    """A single image becomes a one-section Markdown doc with the image
    inlined; the chunker will run VLM on the URL to produce a description."""
    key = f"{storage_prefix}/img-{uuid.uuid4().hex[:8]}.png"
    rel_url = await upload_bytes(key, image_bytes, content_type="image/png")
    if not rel_url:
        return ConvertedDocument(markdown="")

    md = f"# Hình ảnh\n\n![Hình ảnh]({rel_url})\n"
    img = ExtractedImage(
        key=key, url=rel_url, page_number=1,
        order_in_page=1, mime_type="image/png",
    )
    return ConvertedDocument(markdown=md, images=[img], language_hint=language)