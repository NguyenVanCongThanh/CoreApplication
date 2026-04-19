"""
ai-service/app/services/chunker.py
Splits documents (PDF, text, video transcript) into overlapping chunks.
Preserves metadata for Deep Link (page number, video timestamp).
Handles bilingual content (VI/EN).
"""
from __future__ import annotations

import io
import asyncio
import logging
import re
from dataclasses import dataclass
from typing import Callable, Awaitable, Optional

logger = logging.getLogger(__name__)


@dataclass
class DocumentChunk:
    text: str
    index: int
    source_type: str           # 'document' | 'video'
    page_number: int | None = None
    start_time_sec: int | None = None
    end_time_sec: int | None = None
    language: str = "vi"


def sanitize_text(text: str) -> str:
    """
    Remove characters PostgreSQL UTF-8 cannot store.
    - Null bytes (0x00): extracted by PyMuPDF from some PDFs
    - Other non-printable control chars (keep \\n \\r \\t which are fine)
    """
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)


class PDFChunker:
    """
    Chunk PDF files page-by-page, then split large pages into smaller chunks.
    Preserves page_number for Deep Link.
    """

    def __init__(self, chunk_size: int = 500, overlap: int = 50):
        self.chunk_size = chunk_size
        self.overlap = overlap

    def chunk_bytes(self, pdf_bytes: bytes) -> list[DocumentChunk]:
        """Process raw PDF bytes → list of chunks with page metadata."""
        try:
            import pymupdf  # PyMuPDF (faster, better table support)
            return self._chunk_with_pymupdf(pdf_bytes)
        except ImportError:
            logger.warning("PyMuPDF not available, falling back to pypdf")
            try:
                import pypdf
                return self._chunk_with_pypdf(pdf_bytes)
            except ImportError:
                logger.error("No PDF library available")
                return []

    def _chunk_with_pymupdf(self, pdf_bytes: bytes) -> list[DocumentChunk]:
        import pymupdf
        chunks: list[DocumentChunk] = []
        chunk_index = 0

        doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = sanitize_text(page.get_text("text").strip())
            if not text:
                continue

            page_chunks = self._split_text(text)
            for chunk_text in page_chunks:
                chunk_text = sanitize_text(chunk_text.strip())
                if chunk_text:
                    chunks.append(DocumentChunk(
                        text=chunk_text,
                        index=chunk_index,
                        source_type="document",
                        page_number=page_num + 1,
                        language=detect_language(chunk_text),
                    ))
                    chunk_index += 1

        doc.close()
        return chunks

    def _chunk_with_pypdf(self, pdf_bytes: bytes) -> list[DocumentChunk]:
        import pypdf
        chunks: list[DocumentChunk] = []
        chunk_index = 0

        reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
        for page_num, page in enumerate(reader.pages):
            text = sanitize_text((page.extract_text() or "").strip())
            if not text:
                continue

            page_chunks = self._split_text(text)
            for chunk_text in page_chunks:
                chunk_text = sanitize_text(chunk_text.strip())
                if chunk_text:
                    chunks.append(DocumentChunk(
                        text=chunk_text,
                        index=chunk_index,
                        source_type="document",
                        page_number=page_num + 1,
                        language=detect_language(chunk_text),
                    ))
                    chunk_index += 1

        return chunks

    def _split_text(self, text: str) -> list[str]:
        """
        Split text into overlapping chunks at sentence boundaries.
        Prefers sentence breaks over hard character cuts.
        """
        if len(text) <= self.chunk_size:
            return [text]

        # Split into sentences (handles both VI and EN punctuation)
        sentences = re.split(r'(?<=[.!?。！？\n])\s+', text)
        chunks: list[str] = []
        current = ""

        for sentence in sentences:
            if len(current) + len(sentence) <= self.chunk_size:
                current = (current + " " + sentence).strip()
            else:
                if current:
                    chunks.append(current)
                # Start new chunk with overlap from previous
                if current and self.overlap > 0:
                    words = current.split()
                    overlap_text = " ".join(words[-self.overlap // 5:])  # last ~10 words
                    current = (overlap_text + " " + sentence).strip()
                else:
                    current = sentence

        if current:
            chunks.append(current)

        return chunks


class DocxChunker(PDFChunker):
    """Chunk Word (.docx) files."""

    def chunk_bytes(self, docx_bytes: bytes) -> list[DocumentChunk]:
        try:
            from docx import Document
            doc = Document(io.BytesIO(docx_bytes))
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            full_text.append(cell.text.strip())

            text = "\n".join(full_text)
            raw_chunks = self._split_text(text)
            
            return [
                DocumentChunk(
                    text=sanitize_text(c),
                    index=i,
                    source_type="document",
                    page_number=1, # docx doesn't easily provide page numbers
                    language=detect_language(c)
                )
                for i, c in enumerate(raw_chunks)
            ]
        except Exception as e:
            logger.error(f"Error chunking docx: {e}")
            return []


class PptxChunker(PDFChunker):
    """Chunk PowerPoint (.pptx) files."""

    def chunk_bytes(self, pptx_bytes: bytes) -> list[DocumentChunk]:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(pptx_bytes))
            chunks: list[DocumentChunk] = []
            chunk_index = 0

            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Also notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_text.append(notes)
                
                text = "\n".join(slide_text)
                if not text:
                    continue

                slide_chunks = self._split_text(text)
                for c_text in slide_chunks:
                    chunks.append(DocumentChunk(
                        text=sanitize_text(c_text),
                        index=chunk_index,
                        source_type="document",
                        page_number=i + 1,  # slide number as page number
                        language=detect_language(c_text)
                    ))
                    chunk_index += 1
            return chunks
        except Exception as e:
            logger.error(f"Error chunking pptx: {e}")
            return []


class ExcelChunker(PDFChunker):
    """Chunk Excel (.xlsx, .xls) files."""

    def chunk_bytes(self, excel_bytes: bytes) -> list[DocumentChunk]:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(excel_bytes), data_only=True)
            full_text = []
            for sheet in wb.worksheets:
                full_text.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        full_text.append(" | ".join(row_text))
            
            text = "\n".join(full_text)
            raw_chunks = self._split_text(text)
            
            return [
                DocumentChunk(
                    text=sanitize_text(c),
                    index=i,
                    source_type="document",
                    page_number=1,
                    language=detect_language(c)
                )
                for i, c in enumerate(raw_chunks)
            ]
        except Exception as e:
            logger.error(f"Error chunking excel: {e}")
            return []


class VideoTranscriptChunker:
    """
    Chunk video transcripts with timestamp metadata.
    Supports Whisper-style transcripts (SRT or JSON with timestamps).
    Groups segments into ~2-minute chunks for meaningful context.
    """

    def __init__(self, segment_duration_sec: int = 120, overlap_sec: int = 15):
        self.segment_duration = segment_duration_sec
        self.overlap = overlap_sec

    def chunk_whisper_json(self, transcript: dict) -> list[DocumentChunk]:
        """
        Process Whisper JSON transcript:
        {"segments": [{"start": 0.0, "end": 5.2, "text": "..."}, ...]}
        """
        segments = transcript.get("segments", [])
        if not segments:
            return []

        chunks: list[DocumentChunk] = []
        chunk_index = 0
        current_text = ""
        current_start = segments[0]["start"]
        current_end = segments[0]["end"]

        for seg in segments:
            segment_text = sanitize_text(seg.get("text", "").strip())
            seg_start = seg.get("start", 0)
            seg_end = seg.get("end", seg_start)

            # If adding this segment exceeds the chunk duration, flush
            if seg_end - current_start > self.segment_duration and current_text:
                chunks.append(DocumentChunk(
                    text=current_text.strip(),
                    index=chunk_index,
                    source_type="video",
                    start_time_sec=int(current_start),
                    end_time_sec=int(current_end),
                    language=detect_language(current_text),
                ))
                chunk_index += 1
                overlap_start = max(current_start, seg_end - self.overlap)
                current_text = segment_text
                current_start = overlap_start
                current_end = seg_end
            else:
                current_text = (current_text + " " + segment_text).strip()
                current_end = seg_end

        # Flush last chunk
        if current_text.strip():
            chunks.append(DocumentChunk(
                text=current_text.strip(),
                index=chunk_index,
                source_type="video",
                start_time_sec=int(current_start),
                end_time_sec=int(current_end),
                language=detect_language(current_text),
            ))

        return chunks

    def chunk_srt(self, srt_content: str) -> list[DocumentChunk]:
        """Parse SRT format and chunk by duration."""
        blocks = re.split(r"\n\n+", srt_content.strip())
        segments = []
        for block in blocks:
            lines = block.strip().split("\n")
            if len(lines) < 2:
                continue
            for line in lines:
                match = re.match(
                    r"(\d{2}):(\d{2}):(\d{2}),(\d+)\s*-->\s*(\d{2}):(\d{2}):(\d{2}),(\d+)",
                    line,
                )
                if match:
                    h1, m1, s1, ms1, h2, m2, s2, ms2 = match.groups()
                    start = int(h1)*3600 + int(m1)*60 + int(s1)
                    end = int(h2)*3600 + int(m2)*60 + int(s2)
                    text_lines = [l for l in lines if not l.isdigit() and "-->" not in l]
                    text = sanitize_text(" ".join(text_lines).strip())
                    if text:
                        segments.append({"start": start, "end": end, "text": text})
                    break

        return self.chunk_whisper_json({"segments": segments})


def detect_language(text: str) -> str:
    """
    Simple language detection: if >30% of characters are Vietnamese-specific
    diacritics, classify as Vietnamese.
    """
    vi_chars = set("àáảãạăắặằẳẵâấầẩẫậđèéẻẽẹêếềểễệìíỉĩịòóỏõọôốồổỗộơớờởỡợùúủũụưứừửữựỳýỷỹỵ"
                   "ÀÁẢÃẠĂẮẶẰẲẴÂẤẦẨẪẬĐÈÉẺẼẸÊẾỀỂỄỆÌÍỈĨỊÒÓỎÕỌÔỐỒỔỖỘƠỚỜỞỠỢÙÚỦŨỤƯỨỪỬỮỰỲÝỶỸỴ")
    total = len([c for c in text if c.isalpha()])
    if total == 0:
        return "vi"
    vi_count = sum(1 for c in text if c in vi_chars)
    return "vi" if vi_count / total > 0.05 else "en"


def format_timestamp(seconds: int) -> str:
    """Convert seconds to MM:SS or HH:MM:SS string."""
    h = seconds // 3600
    m = (seconds % 3600) // 60
    s = seconds % 60
    if h > 0:
        return f"{h:02d}:{m:02d}:{s:02d}"
    return f"{m:02d}:{s:02d}"

class MarkdownChunker(PDFChunker):   # type: ignore[name-defined]
    """
    Semantic chunker for Markdown text (TEXT content type).
 
    Pipeline (async path):
      1. Protect code blocks (replaced with placeholders during processing)
      2. Replace image references with VLM descriptions (async, concurrent)
      3. Add semantic context prefix to tables
      4. Split at heading boundaries, building heading breadcrumbs
      5. Sub-chunk oversized sections at paragraph / sentence boundaries
      6. Restore code blocks
 
    Sync path (chunk_bytes): images replaced with alt text only — suitable
    for Celery workers that don't have an async event loop.
    """
 
    # ── Regex patterns ────────────────────────────────────────────────────────
 
    _RE_IMAGE = re.compile(r"!\[([^\]]*)\]\(([^)\s]+)[^)]*\)")
    _RE_HEADING = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)
    _RE_TABLE_BLOCK = re.compile(
        r"(\|[^\n]+\|\n)(\|[-:| ]+\|\n)((?:\|[^\n]+\|\n)*)",
        re.MULTILINE,
    )
    _RE_CODE_BLOCK = re.compile(r"```([^\n]*)\n([\s\S]*?)```", re.MULTILINE)
    _RE_BLANK_LINES = re.compile(r"\n{3,}")
 
    _MIN_CHUNK_LEN = 30
 
    # ── Public interface ──────────────────────────────────────────────────────
 
    async def chunk_async(
        self,
        markdown_text: str,
        image_describer: Optional[Callable[[str, str], Awaitable[str]]] = None,
    ) -> list:  # list[DocumentChunk]
        """
        Full async pipeline — use this from FastAPI / auto_index_service.
        image_describer(url, alt_text) -> description string.
        """
        text = markdown_text
 
        # Step 1: protect code blocks
        text, code_blocks = self._extract_code_blocks(text)
 
        # Step 2: replace images with VLM descriptions (concurrent)
        text = await self._process_images(text, image_describer)
 
        # Step 3: table semantic prefix
        text = self._add_table_context(text)
 
        # Step 4–5: heading-based semantic split
        chunks = self._heading_split(text)
 
        # Step 6: restore code blocks inside chunk texts
        chunks = self._restore_code_blocks(chunks, code_blocks)
 
        return chunks
 
    def chunk_bytes(self, text_bytes: bytes) -> list:  # list[DocumentChunk]
        """
        Sync fallback — images get alt text, no VLM.
        Used by Celery process_document_task.
        """
        text = text_bytes.decode("utf-8", errors="replace")
        text, code_blocks = self._extract_code_blocks(text)
        text = self._replace_images_sync(text)
        text = self._add_table_context(text)
        chunks = self._heading_split(text)
        chunks = self._restore_code_blocks(chunks, code_blocks)
        return chunks
 
    # ── Step 1: code block protection ────────────────────────────────────────
 
    def _extract_code_blocks(self, text: str) -> tuple[str, dict[str, str]]:
        """Replace code blocks with placeholders; return (text, placeholder_map)."""
        placeholders: dict[str, str] = {}
 
        def _replace(m: re.Match) -> str:
            lang = m.group(1).strip() or "code"
            code = m.group(2)
            key = f"\x00CODE_{len(placeholders)}\x00"
            # Store with a helpful label so it reads naturally in the chunk
            placeholders[key] = f"[Đoạn code {lang}]:\n```{lang}\n{code}```"
            return key
 
        return self._RE_CODE_BLOCK.sub(_replace, text), placeholders
 
    def _restore_code_blocks(self, chunks: list, placeholders: dict[str, str]) -> list:
        restored = []
        for chunk in chunks:
            text = chunk.text
            for key, value in placeholders.items():
                text = text.replace(key, value)
            chunk.text = sanitize_text(text)  # type: ignore[name-defined]
            restored.append(chunk)
        return restored
 
    # ── Step 2: image processing ──────────────────────────────────────────────
 
    async def _process_images(
        self,
        text: str,
        image_describer: Optional[Callable],
    ) -> str:
        """Replace all ![alt](url) with VLM descriptions (concurrent)."""
        matches = list(self._RE_IMAGE.finditer(text))
        if not matches:
            return text
 
        if image_describer is None:
            return self._replace_images_sync(text)
 
        # Call VLM concurrently for all images
        tasks = [
            image_describer(m.group(2), m.group(1))  # (url, alt_text)
            for m in matches
        ]
        results = await asyncio.gather(*tasks, return_exceptions=True)
 
        # Rebuild text by replacing matches in reverse order (preserve indices)
        result = text
        for match, desc in zip(reversed(matches), reversed(results)):
            if isinstance(desc, Exception) or not desc:
                alt = match.group(1) or "hình ảnh"
                replacement = f"\n[Hình ảnh: {alt}]\n"
            else:
                replacement = f"\n[Mô tả hình ảnh: {desc}]\n"
            result = result[: match.start()] + replacement + result[match.end() :]
 
        return result
 
    def _replace_images_sync(self, text: str) -> str:
        """Sync fallback: replace images with alt text."""
        def _sub(m: re.Match) -> str:
            alt = m.group(1).strip()
            return f"\n[Hình ảnh: {alt or 'không có mô tả'}]\n"
 
        return self._RE_IMAGE.sub(_sub, text)
 
    # ── Step 3: table semantic context ────────────────────────────────────────
 
    def _add_table_context(self, text: str) -> str:
        """
        Prepend a natural-language header summary to each markdown table.
 
        Before:
          | Tên | Tuổi | Điểm |
          |-----|------|------|
          | An  | 20   | 8.5  |
 
        After:
          [Bảng dữ liệu - Các cột: Tên, Tuổi, Điểm]
          | Tên | Tuổi | Điểm |
          ...
        """
        def _replace(m: re.Match) -> str:
            header_row = m.group(1)
            # Parse column names from | col1 | col2 | col3 |
            cols = [c.strip() for c in header_row.strip().strip("|").split("|") if c.strip()]
            if cols:
                summary = f"[Bảng dữ liệu — Các cột: {', '.join(cols)}]\n"
            else:
                summary = "[Bảng dữ liệu]\n"
            return summary + m.group(0)
 
        return self._RE_TABLE_BLOCK.sub(_replace, text)
 
    # ── Steps 4–5: heading-based semantic split ────────────────────────────────
 
    def _heading_split(self, text: str) -> list:  # list[DocumentChunk]
        """
        Split markdown at heading boundaries.
        Each resulting chunk is prefixed with its heading breadcrumb so the
        embedding knows which section this passage belongs to even without
        surrounding context.
        """
        # Normalize blank lines
        text = self._RE_BLANK_LINES.sub("\n\n", text)
 
        # Split text at heading lines, keeping the heading with its section
        # (?=...) lookahead preserves the heading line in the next split token
        parts = re.split(r"(?=^#{1,6}\s)", text, flags=re.MULTILINE)
 
        chunks = []
        chunk_idx = 0
        heading_stack: list[tuple[int, str]] = []  # [(level, title), ...]
 
        for part in parts:
            if not part.strip():
                continue
 
            # Extract leading heading (if any)
            heading_match = self._RE_HEADING.match(part.strip())
            if heading_match:
                level = len(heading_match.group(1))
                title = heading_match.group(2).strip()
                # Pop headings of equal or deeper level
                heading_stack = [(l, t) for l, t in heading_stack if l < level]
                heading_stack.append((level, title))
 
                # Body text = everything after the first heading line
                lines = part.split("\n", 1)
                body = lines[1].strip() if len(lines) > 1 else ""
            else:
                body = part.strip()
 
            if not body:
                continue
 
            # Build breadcrumb
            breadcrumb = " > ".join(t for _, t in heading_stack)
            prefix = f"[{breadcrumb}]\n" if breadcrumb else ""
 
            # Combine prefix + body
            full_text = prefix + body
 
            # Sub-chunk if section is too large
            sub_texts = self._split_text(full_text)
            for sub in sub_texts:
                sub = sanitize_text(sub.strip())  # type: ignore[name-defined]
                if len(sub) < self._MIN_CHUNK_LEN:
                    continue
 
                # Each sub-chunk still carries the breadcrumb if it was stripped
                if prefix and not sub.startswith("["):
                    sub = prefix.rstrip("\n") + "\n" + sub
 
                chunks.append(
                    DocumentChunk(  # type: ignore[name-defined]
                        text=sub,
                        index=chunk_idx,
                        source_type="document",
                        page_number=None,
                        language=detect_language(sub),  # type: ignore[name-defined]
                    )
                )
                chunk_idx += 1
 
        return chunks
 
 
# ─────────────────────────────────────────────────────────────────────────────
 
 
class ImageChunker:
    """
    Handles standalone IMAGE content type.
 
    Async path: calls VLM to generate a rich description.
    Sync path:  returns a placeholder chunk (re-index later with VLM).
 
    The description is stored as a single DocumentChunk; the auto_index_service
    then runs node extraction on it just like any other text.
    """
 
    async def chunk_async(
        self,
        image_bytes: bytes,
        mime_type: str = "image/jpeg",
        language: str = "vi",
    ) -> list:  # list[DocumentChunk]
        from app.core.vlm import describe_image_bytes
 
        description = await describe_image_bytes(image_bytes, language=language, mime_type=mime_type)
 
        return [
            DocumentChunk(  # type: ignore[name-defined]
                text=description,
                index=0,
                source_type="document",
                page_number=1,
                language=detect_language(description),  # type: ignore[name-defined]
            )
        ]
 
    def chunk_bytes(self, image_bytes: bytes) -> list:  # list[DocumentChunk]
        """Sync fallback — used by Celery. Returns placeholder."""
        return [
            DocumentChunk(  # type: ignore[name-defined]
                text="[Hình ảnh chưa được mô tả — cần chạy lại pipeline với VLM]",
                index=0,
                source_type="document",
                page_number=1,
                language="vi",
            )
        ]
 